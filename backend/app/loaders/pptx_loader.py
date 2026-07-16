"""PPTX 加载器：提取文本、表格和图片。"""

from __future__ import annotations

from pathlib import Path

from app.loaders.base import BaseLoader, ImageInfo, LoadedDocument, LoadedPage


class PptxLoader(BaseLoader):
    extensions = (".pptx",)

    def load(self, path: str | Path) -> LoadedDocument:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        p = Path(path)
        prs = Presentation(str(p))

        pages: list[LoadedPage] = []
        img_idx = 0

        for i, slide in enumerate(prs.slides, start=1):
            buf: list[str] = []
            images: list[ImageInfo] = []
            title = None

            for shape in slide.shapes:
                # 文本
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if not text:
                            continue
                        if title is None and len(text) < 80:
                            title = text
                        buf.append(text)

                # 表格
                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            rows.append(" | ".join(cells))
                    if rows:
                        buf.append("\n".join(rows))

                # 图片
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        img_bytes = image.blob
                        mime = "image/png"
                        if "jpeg" in image.content_type or "jpg" in image.content_type:
                            mime = "image/jpeg"
                        elif "png" in image.content_type:
                            mime = "image/png"
                        elif image.ext == "jpg":
                            mime = "image/jpeg"

                        img_idx += 1
                        img_info = ImageInfo(
                            image_bytes=img_bytes,
                            mime_type=mime,
                            position=len("\n".join(buf)),
                        )
                        images.append(img_info)
                        buf.append(f"[图片:{img_idx}]")
                    except Exception:
                        continue

            pages.append(
                LoadedPage(
                    text="\n".join(buf),
                    page=i,
                    section=title,
                    metadata={"page": i, "title": title},
                    images=images,
                )
            )

        full = "\n\n".join(p.text for p in pages if p.text)
        total_images = sum(len(p.images) for p in pages)

        return LoadedDocument(
            text=full,
            pages=pages,
            metadata={"slide_count": len(pages), "image_count": total_images},
            source=str(p),
            ext=p.suffix.lower(),
        )